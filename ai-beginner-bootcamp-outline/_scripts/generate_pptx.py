import os
import json
import httpx
from anthropic import Anthropic
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Load environment variables
load_dotenv()

api_key = os.getenv("AGENT_ROUTER_API_KEY")
if not api_key:
    print("Error: AGENT_ROUTER_API_KEY not found in .env")
    exit(1)

# Configure the Anthropic client to point to Agent Router
client = Anthropic(
    api_key=api_key,
    base_url="https://agentrouter.org",
    http_client=httpx.Client(
        headers={
            "User-Agent": "cline/2.0.0",
            "Accept": "application/json"
        },
        timeout=60.0
    )
)

# AgenticLabs Brand Colors
BG_COLOR = RGBColor(11, 17, 32)      # #0B1120
ACCENT_COLOR = RGBColor(0, 229, 255) # #00E5FF
TEXT_COLOR = RGBColor(255, 255, 255) # White

def apply_slide_styling(slide, title, prs):
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # Set title properties if the slide has a title shape
    if slide.shapes.title:
        title_shape = slide.shapes.title
        title_shape.text = title
        if title_shape.text_frame.paragraphs:
            p = title_shape.text_frame.paragraphs[0]
            p.font.name = "Arial"
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = ACCENT_COLOR

def generate_slide_content(day_number, topic):
    prompt = f"""
    You are an expert curriculum designer for AgenticLabs.ng.
    We are building a 5-Day "Second Brain" AI Bootcamp.
    
    Day {day_number} Topic: {topic}
    
    Generate the content for 6 to 8 presentation slides for this day. 
    The tone should be encouraging, non-technical at first but building to clear concepts.
    
    Respond STRICTLY with a JSON object in this exact format:
    {{
        "slides": [
            {{
                "title": "Slide Title",
                "bullets": ["Bullet 1", "Bullet 2", "Bullet 3"],
                "presenter_notes": "What the speaker should say for this slide."
            }}
        ]
    }}
    Do not include any other text or markdown formatting around the JSON.
    """
    
    try:
        response = client.messages.create(
            model="claude-opus-5",
            max_tokens=8000,
            messages=[{"role": "user", "content": prompt}]
        )
        
        content = ""
        for block in response.content:
            if getattr(block, "type", "") == "text":
                content = block.text
            elif hasattr(block, "text"):
                content = block.text
                
        content = content.strip()
        if content.startswith("```json"):
            content = content[7:-3].strip()
        elif content.startswith("```"):
            content = content[3:-3].strip()
            
        data = json.loads(content)
        return data.get("slides", [])
    except Exception as e:
        print(f"Error generating content for Day {day_number}: {e}")
        return []

def create_pptx_for_day(day_number, topic):
    print(f"Generating content for Day {day_number}...")
    slides_data = generate_slide_content(day_number, topic)
    
    if not slides_data:
        print(f"Failed to generate slides for Day {day_number}")
        return

    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    apply_slide_styling(slide, f"Day {day_number}: {topic}", prs)
    
    subtitle = slide.placeholders[1]
    subtitle.text = "AgenticLabs.ng AI Bootcamp"
    if subtitle.text_frame.paragraphs:
        p = subtitle.text_frame.paragraphs[0]
        p.font.color.rgb = TEXT_COLOR
        p.font.size = Pt(24)

    # Content Slides
    bullet_slide_layout = prs.slide_layouts[1]
    for s_data in slides_data:
        slide = prs.slides.add_slide(bullet_slide_layout)
        apply_slide_styling(slide, s_data["title"], prs)
        
        # Add bullets
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for bullet in s_data.get("bullets", []):
            p = tf.add_paragraph()
            p.text = bullet
            p.font.name = "Arial"
            p.font.size = Pt(24)
            p.font.color.rgb = TEXT_COLOR
            p.level = 0
            
        # Add presenter notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = s_data.get("presenter_notes", "")

    # Save
    out_dir = f"../day_{day_number}"
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, f"Day_{day_number}_Slides.pptx")
    prs.save(out_path)
    print(f"Saved {out_path}")

def main():
    days = [
        (1, "Hire the Brain (Intro & Setup)"),
        (2, "Give it Hands (Tools & Search)"),
        (3, "Give it Memory (Sessions & Context)"),
        (4, "Quality Control & Guardrails"),
        (5, "Multi-Agent Teams & Launch")
    ]
    
    for day_num, topic in days:
        create_pptx_for_day(day_num, topic)

if __name__ == "__main__":
    main()
